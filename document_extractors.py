import csv
import io
import re
from dataclasses import dataclass, field
from typing import List, Optional


DEVANAGARI_PATTERN = re.compile(r"[\u0900-\u097F]")


@dataclass
class ExtractedSection:
    text: str
    page: Optional[int] = None
    section: str = ""


@dataclass
class ExtractedDocument:
    title: str
    sections: List[ExtractedSection] = field(default_factory=list)
    file_name: str = ""
    source_type: str = ""

    @property
    def full_text(self):
        parts = []
        for item in self.sections:
            prefix = ""
            if item.section:
                prefix = f"## {item.section}\n"
            elif item.page:
                prefix = f"[Page {item.page}]\n"
            parts.append(prefix + item.text)
        return "\n\n".join(parts)


def detect_language(text):
    sample = (text or "")[:4000]
    if not sample.strip():
        return "English"

    devanagari = len(DEVANAGARI_PATTERN.findall(sample))
    if devanagari == 0:
        return "English"
    if devanagari > len(sample) * 0.08:
        return "Hindi"
    return "Mixed"


def extract_pdf(uploaded_file, max_pages=80):
    from pypdf import PdfReader

    reader = PdfReader(uploaded_file)
    sections = []
    for page_number, page in enumerate(reader.pages[:max_pages], start=1):
        text = (page.extract_text() or "").strip()
        if text:
            sections.append(ExtractedSection(text=text, page=page_number))

    name = getattr(uploaded_file, "name", "document.pdf")
    return ExtractedDocument(title=name, sections=sections, file_name=name, source_type="pdf")


def extract_docx(uploaded_file):
    from docx import Document

    document = Document(uploaded_file)
    sections = []
    current_heading = ""
    buffer = []

    def flush_buffer():
        nonlocal buffer, current_heading
        if buffer:
            sections.append(
                ExtractedSection(
                    text="\n".join(buffer),
                    section=current_heading,
                )
            )
            buffer = []

    for paragraph in document.paragraphs:
        text = " ".join(paragraph.text.split())
        if not text:
            continue

        style_name = (paragraph.style.name or "").lower()
        if "heading" in style_name:
            flush_buffer()
            current_heading = text
            continue

        buffer.append(text)

    flush_buffer()

    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [" ".join(cell.text.split()) for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sections.append(
                ExtractedSection(
                    text="\n".join(rows),
                    section=current_heading or "Table",
                )
            )

    name = getattr(uploaded_file, "name", "document.docx")
    return ExtractedDocument(title=name, sections=sections, file_name=name, source_type="docx")


def extract_txt(uploaded_file):
    raw = uploaded_file.read()
    if isinstance(raw, bytes):
        text = raw.decode("utf-8", errors="replace")
    else:
        text = str(raw)

    name = getattr(uploaded_file, "name", "document.txt")
    sections = [ExtractedSection(text=text.strip())] if text.strip() else []
    return ExtractedDocument(title=name, sections=sections, file_name=name, source_type="txt")


def extract_csv(uploaded_file):
    raw = uploaded_file.read()
    if isinstance(raw, bytes):
        content = raw.decode("utf-8", errors="replace")
    else:
        content = str(raw)

    reader = csv.reader(io.StringIO(content))
    rows = []
    for row in reader:
        if any(cell.strip() for cell in row):
            rows.append(" | ".join(cell.strip() for cell in row))

    name = getattr(uploaded_file, "name", "document.csv")
    text = "\n".join(rows)
    sections = [ExtractedSection(text=text, section="CSV Data")] if text else []
    return ExtractedDocument(title=name, sections=sections, file_name=name, source_type="csv")


def extract_xlsx(uploaded_file):
    from openpyxl import load_workbook

    workbook = load_workbook(uploaded_file, read_only=True, data_only=True)
    sections = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sections.append(ExtractedSection(text="\n".join(rows), section=f"Sheet: {sheet.title}"))

    name = getattr(uploaded_file, "name", "document.xlsx")
    return ExtractedDocument(title=name, sections=sections, file_name=name, source_type="xlsx")


def extract_pptx(uploaded_file):
    from pptx import Presentation

    presentation = Presentation(uploaded_file)
    sections = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = " ".join(shape.text.split())
                if text:
                    lines.append(text)

        if lines:
            sections.append(
                ExtractedSection(
                    text="\n".join(lines),
                    page=slide_number,
                    section=f"Slide {slide_number}",
                )
            )

    name = getattr(uploaded_file, "name", "document.pptx")
    return ExtractedDocument(title=name, sections=sections, file_name=name, source_type="pptx")


SUPPORTED_FILE_TYPES = {
    "pdf": (["pdf"], extract_pdf),
    "docx": (["docx"], extract_docx),
    "txt": (["txt"], extract_txt),
    "csv": (["csv"], extract_csv),
    "xlsx": (["xlsx"], extract_xlsx),
    "pptx": (["pptx"], extract_pptx),
}


def extract_document(source_type, uploaded_file):
    source_type = (source_type or "").lower().strip()
    if source_type not in SUPPORTED_FILE_TYPES:
        raise ValueError(f"Unsupported file type: {source_type}")

    _, extractor = SUPPORTED_FILE_TYPES[source_type]
    return extractor(uploaded_file)
